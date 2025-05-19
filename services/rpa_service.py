import subprocess
import os
import shutil
import tempfile
from selenium import webdriver
from selenium.webdriver.chrome.options import Options
from selenium.webdriver.chrome.service import Service
from selenium.webdriver.common.by import By
from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC
import pandas as pd
import json
import datetime
import logging
import io
import time
from pymongo import MongoClient
from db.database import database
from webdriver_manager.chrome import ChromeDriverManager
import sys


logging.basicConfig(level=logging.DEBUG, format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

class RPAServiceError(Exception):
    """Excepción personalizada para errores en el servicio RPA."""
    pass

def check_rpa_dependencies():
    """Verifica que las dependencias para RPA estén instaladas."""
    try:
        # Verificar Python
        logger.info("Verificando versión de Python...")
        python_version = subprocess.check_output(["python", "--version"]).decode().strip()
        logger.info(f"Versión de Python: {python_version}")
        
        # Verificar Chrome usando webdriver_manager
        try:
            chrome_path = "Gestionado automáticamente por webdriver_manager"
            chromedriver_path = ChromeDriverManager().install()
            logger.info(f"ChromeDriver instalado en: {chromedriver_path}")
        except Exception as e:
            logger.warning(f"Error al verificar ChromeDriver: {str(e)}")
            chrome_path = None
            chromedriver_path = None
        
        # Verificar Selenium
        try:
            import selenium
            logger.info(f"Selenium versión: {selenium.__version__}")
        except ImportError:
            logger.error("Selenium no está instalado")
            raise RPAServiceError("Selenium no está instalado")
        
        # Verificar pywin32 para automatización de Office
        office_automation_available = False
        try:
            import win32com.client
            office_automation_available = True
            logger.info("pywin32 disponible para automatización de Office")
        except ImportError:
            logger.warning("pywin32 no está instalado. La automatización de Office no estará disponible.")
        
        return {
            "python_version": python_version,
            "chrome_path": chrome_path or "No encontrado",
            "chromedriver_path": chromedriver_path or "No encontrado",
            "selenium_version": selenium.__version__,
            "office_automation_available": office_automation_available,
            "status": "Todas las dependencias disponibles" if chrome_path and chromedriver_path else "Faltan algunas dependencias"
        }
    except Exception as e:
        logger.error(f"Error al verificar dependencias de RPA: {str(e)}")
        raise RPAServiceError(f"Error al verificar dependencias: {str(e)}")

def scrape_data_automated(url, selector_type="css", selector_value=None, wait_time=10, process_fields=None):
    """
    Realiza scraping de datos de forma automatizada utilizando Selenium.
    
    Args:
        url (str): URL a scrapear
        selector_type (str): Tipo de selector (css, xpath, id, class)
        selector_value (str): Valor del selector
        wait_time (int): Tiempo máximo de espera en segundos
        process_fields (list): Lista de campos a procesar
        
    Returns:
        tuple[io.BytesIO, str]: Buffer de Excel y nombre de archivo
    """
    logger.info(f"Iniciando scraping automatizado de: {url}")
    
    options = Options()
    options.add_argument("--headless")
    options.add_argument("--no-sandbox")
    options.add_argument("--disable-dev-shm-usage")
    options.add_argument("--disable-gpu")
    options.add_argument("--window-size=1920,1080")
    
    try:
        # Usar webdriver_manager para descargar y gestionar ChromeDriver automáticamente
        service = Service(ChromeDriverManager().install())
        driver = webdriver.Chrome(service=service, options=options)
        driver.get(url)
        
        # Esperar a que la página cargue completamente
        driver.implicitly_wait(wait_time)
        
        # Esperar a que el elemento selector aparezca
        if selector_value:
            by_type = {
                "css": By.CSS_SELECTOR,
                "xpath": By.XPATH,
                "id": By.ID,
                "class": By.CLASS_NAME
            }.get(selector_type.lower(), By.CSS_SELECTOR)
            
            WebDriverWait(driver, wait_time).until(
                EC.presence_of_element_located((by_type, selector_value))
            )
        
        # Extraer datos
        data = []
        
        if process_fields:
            # Procesamiento personalizado según los campos especificados
            for field in process_fields:
                if field.get('selector'):
                    try:
                        elements = driver.find_elements(
                            by={
                                "css": By.CSS_SELECTOR,
                                "xpath": By.XPATH,
                                "id": By.ID,
                                "class": By.CLASS_NAME
                            }.get(field.get('selector_type', 'css'), By.CSS_SELECTOR),
                            value=field['selector']
                        )
                        
                        field_data = []
                        for element in elements:
                            if field.get('attribute'):
                                value = element.get_attribute(field['attribute'])
                            else:
                                value = element.text
                            field_data.append(value)
                        
                        field['values'] = field_data
                    except Exception as e:
                        logger.warning(f"Error al extraer campo {field.get('name')}: {str(e)}")
                        field['values'] = []
            
            # Convertir a estructura tabular
            if process_fields and all('values' in field for field in process_fields):
                max_length = max(len(field['values']) for field in process_fields)
                table_data = []
                
                for i in range(max_length):
                    row = {}
                    for field in process_fields:
                        values = field['values']
                        row[field.get('name', field.get('selector', 'unknown'))] = values[i] if i < len(values) else None
                    table_data.append(row)
                
                data = table_data
        else:
            # Extracción automática de datos de tablas
            tables = driver.find_elements(By.TAG_NAME, 'table')
            
            if tables:
                logger.info(f"Se encontraron {len(tables)} tablas en la página")
                
                for i, table in enumerate(tables):
                    try:
                        rows = table.find_elements(By.TAG_NAME, 'tr')
                        table_data = []
                        
                        # Extraer encabezados
                        headers = []
                        header_row = rows[0] if rows else None
                        if header_row:
                            header_cells = header_row.find_elements(By.TAG_NAME, 'th')
                            if not header_cells:
                                header_cells = header_row.find_elements(By.TAG_NAME, 'td')
                            headers = [cell.text.strip() for cell in header_cells]
                        
                        # Extraer filas de datos
                        for row in rows[1:]:
                            cells = row.find_elements(By.TAG_NAME, 'td')
                            if cells:
                                row_data = {headers[j] if j < len(headers) else f"Column_{j}": cells[j].text.strip() 
                                           for j in range(len(cells))}
                                table_data.append(row_data)
                        
                        data.extend(table_data)
                    except Exception as e:
                        logger.warning(f"Error al procesar tabla {i}: {str(e)}")
            else:
                # Si no hay tablas, intentar extraer datos estructurados
                logger.info("No se encontraron tablas, intentando extraer datos estructurados")
                
                elements = driver.find_elements(By.CSS_SELECTOR, 'div[class*="item"], div[class*="card"], article')
                
                if elements:
                    logger.info(f"Se encontraron {len(elements)} elementos estructurados")
                    
                    for element in elements:
                        try:
                            title_el = element.find_elements(By.CSS_SELECTOR, 'h1, h2, h3, h4, .title')
                            desc_el = element.find_elements(By.CSS_SELECTOR, 'p, .description')
                            
                            element_data = {
                                "title": title_el[0].text.strip() if title_el else "",
                                "description": desc_el[0].text.strip() if desc_el else "",
                                "html": element.get_attribute('outerHTML')
                            }
                            
                            data.append(element_data)
                        except Exception as e:
                            logger.warning(f"Error al procesar elemento: {str(e)}")
                else:
                    # Si no hay elementos estructurados, extraer todo el texto de la página
                    logger.info("No se encontraron elementos estructurados, extrayendo todo el texto")
                    
                    body_text = driver.find_element(By.TAG_NAME, 'body').text
                    data = [{"full_text": body_text}]
        
        # Crear DataFrame
        if isinstance(data, list) and data:
            df = pd.DataFrame(data)
        else:
            # Si no se encontraron datos estructurados, usar todo el HTML
            logger.warning("No se encontraron datos estructurados, usando todo el HTML")
            df = pd.DataFrame([{"html": driver.page_source}])
        
        # Cerrar el navegador
        driver.quit()
        
        # Crear archivo Excel
        output = io.BytesIO()
        df.to_excel(output, index=False)
        output.seek(0)
        
        # Nombre del archivo
        timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"rpa_scraping_{timestamp}.xlsx"
        
        return output, filename
    
    except Exception as e:
        logger.error(f"Error en el scraping automatizado: {str(e)}")
        # Asegurarse de cerrar el navegador en caso de error
        try:
            driver.quit()
        except:
            pass
        raise RPAServiceError(f"Error en el scraping automatizado: {str(e)}")

def automate_data_entry(collection_name, data_source, field_mappings=None):
    """
    Automatiza la entrada de datos desde un origen a una colección MongoDB.
    
    Args:
        collection_name (str): Nombre de la colección donde insertar datos
        data_source (str): Fuente de datos (URL o archivo)
        field_mappings (dict): Mapeo de campos de origen a destino
        
    Returns:
        dict: Resultado de la operación
    """
    logger.info(f"Iniciando automatización de entrada de datos a colección: {collection_name}")
    
    try:
        # Determinar el tipo de origen de datos
        is_url = data_source.startswith('http://') or data_source.startswith('https://')
        
        # Cargar datos
        if is_url:
            # Si es una URL, hacer scraping
            logger.info(f"Origen de datos es una URL: {data_source}")
            
            buffer, _ = scrape_data_automated(data_source)
            df = pd.read_excel(buffer)
        else:
            # Si es un archivo local
            logger.info(f"Origen de datos es un archivo: {data_source}")
            
            if data_source.endswith('.csv'):
                df = pd.read_csv(data_source)
            elif data_source.endswith('.xlsx') or data_source.endswith('.xls'):
                df = pd.read_excel(data_source)
            else:
                raise RPAServiceError(f"Formato de archivo no soportado: {data_source}")
        
        # Verificar que tenemos datos
        if df.empty:
            raise RPAServiceError("No se encontraron datos en el origen")
        
        logger.info(f"Se cargaron {len(df)} registros de datos")
        
        # Aplicar mapeo de campos si se proporciona
        if field_mappings:
            logger.info(f"Aplicando mapeo de campos: {field_mappings}")
            
            # Renombrar columnas según el mapeo
            rename_dict = {}
            for dest_field, source_field in field_mappings.items():
                if source_field in df.columns:
                    rename_dict[source_field] = dest_field
            
            if rename_dict:
                df = df.rename(columns=rename_dict)
            
            # Seleccionar solo las columnas mapeadas
            df = df[[dest_field for dest_field in field_mappings.values() 
                    if dest_field in df.columns]]
        
        # Convertir DataFrame a lista de diccionarios
        records = df.to_dict('records')
        
        # Insertar en MongoDB
        collection = database[collection_name]
        result = collection.insert_many(records)
        
        logger.info(f"Se insertaron {len(result.inserted_ids)} registros en la colección {collection_name}")
        
        return {
            "success": True,
            "collection": collection_name,
            "records_processed": len(df),
            "records_inserted": len(result.inserted_ids),
            "inserted_ids": [str(id) for id in result.inserted_ids]
        }
    
    except Exception as e:
        logger.error(f"Error en la automatización de entrada de datos: {str(e)}")
        raise RPAServiceError(f"Error en la automatización de entrada de datos: {str(e)}")

def schedule_data_sync(source_collection, target_collection, sync_interval_minutes=60, field_mappings=None):
    """
    Configura una sincronización programada entre colecciones o fuentes de datos.
    
    Args:
        source_collection (str): Nombre de la colección o URL fuente
        target_collection (str): Nombre de la colección destino
        sync_interval_minutes (int): Intervalo de sincronización en minutos
        field_mappings (dict): Mapeo de campos
        
    Returns:
        dict: Información de la tarea programada
    """
    logger.info(f"Configurando sincronización: {source_collection} -> {target_collection} cada {sync_interval_minutes} minutos")
    
    try:
        # Crear un archivo temporal para almacenar la configuración
        config_dir = os.path.join(tempfile.gettempdir(), "rpa_sync_configs")
        os.makedirs(config_dir, exist_ok=True)
        
        config = {
            "source": source_collection,
            "target": target_collection,
            "interval_minutes": sync_interval_minutes,
            "field_mappings": field_mappings,
            "created_at": datetime.datetime.now().isoformat(),
            "last_sync": None,
            "next_sync": (datetime.datetime.now() + datetime.timedelta(minutes=sync_interval_minutes)).isoformat()
        }
        
        # Generar ID único para la tarea
        task_id = f"sync_{int(time.time())}"
        config_file = os.path.join(config_dir, f"{task_id}.json")
        
        # Guardar configuración
        with open(config_file, 'w') as f:
            json.dump(config, f)
        
        logger.info(f"Configuración de sincronización guardada en: {config_file}")
        
        # Si source_collection es una URL, probar que podemos acceder
        if source_collection.startswith('http://') or source_collection.startswith('https://'):
            try:
                import requests
                response = requests.head(source_collection, timeout=5)
                response.raise_for_status()
            except Exception as e:
                logger.warning(f"La URL de origen no es accesible: {str(e)}")
                
        # Almacenar la configuración en MongoDB para persistencia
        sync_collection = database.get_collection("rpa_sync_tasks")
        sync_collection.insert_one({
            "task_id": task_id,
            "config": config,
            "status": "scheduled"
        })
        
        return {
            "task_id": task_id,
            "source": source_collection,
            "target": target_collection,
            "interval_minutes": sync_interval_minutes,
            "next_sync": config["next_sync"],
            "status": "scheduled"
        }
    
    except Exception as e:
        logger.error(f"Error al programar sincronización: {str(e)}")
        raise RPAServiceError(f"Error al programar sincronización: {str(e)}")

def execute_rpa_script(script_content, script_type="python", params=None):
    """
    Ejecuta un script RPA personalizado.
    
    Args:
        script_content (str): Contenido del script
        script_type (str): Tipo de script (python, javascript)
        params (dict): Parámetros para el script
        
    Returns:
        dict: Resultado de la ejecución
    """
    logger.info(f"Ejecutando script RPA de tipo: {script_type}")
    
    try:
        # Crear directorio temporal para el script
        script_dir = os.path.join(tempfile.gettempdir(), "rpa_scripts")
        os.makedirs(script_dir, exist_ok=True)
        
        # Crear archivo temporal para el script
        if script_type.lower() == "python":
            script_file = os.path.join(script_dir, f"rpa_script_{int(time.time())}.py")
        elif script_type.lower() == "javascript":
            script_file = os.path.join(script_dir, f"rpa_script_{int(time.time())}.js")
        else:
            raise RPAServiceError(f"Tipo de script no soportado: {script_type}")
        
        # Guardar script en archivo temporal con codificación UTF-8 explícita
        with open(script_file, 'w', encoding='utf-8') as f:
            f.write(script_content)
        
        logger.info(f"Script guardado en: {script_file}")
        
        # Preparar parámetros como argumentos de línea de comandos
        cmd_args = []
        if params:
            for key, value in params.items():
                if isinstance(value, (str, int, float, bool)):
                    cmd_args.append(f"--{key}={value}")
        
        # Ejecutar script
        if script_type.lower() == "python":
            cmd = ["python", script_file] + cmd_args
        elif script_type.lower() == "javascript":
            cmd = ["node", script_file] + cmd_args
        
        logger.info(f"Ejecutando comando: {' '.join(cmd)}")
        
        # Crear archivo .bat para ejecutar el script (solución específica para Windows)
        if os.name == 'nt' and script_type.lower() == "python":
            bat_file = os.path.join(script_dir, f"rpa_script_{int(time.time())}.bat")
            with open(bat_file, 'w', encoding='cp1252') as f:
                f.write(f'@echo off\r\nchcp 65001\r\npython "{script_file}" {" ".join(cmd_args)}\r\n')
            
            cmd = [bat_file]
        
        # Iniciar proceso con límite de tiempo (30 minutos por defecto)
        # Usar CP1252 (Windows-1252) para la codificación de salida en Windows
        # y errors='replace' para manejar errores de decodificación
        process = subprocess.Popen(
            cmd, 
            stdout=subprocess.PIPE, 
            stderr=subprocess.PIPE,
            text=True,
            encoding='cp1252' if os.name == 'nt' else 'utf-8',  # Usar cp1252 en Windows
            errors='replace'  # Reemplazar caracteres que no se pueden decodificar
        )
        
        try:
            stdout, stderr = process.communicate(timeout=1800)  # 30 minutos
        except subprocess.TimeoutExpired:
            process.kill()
            stdout, stderr = process.communicate()
            logger.error("La ejecución del script excedió el tiempo límite")
            raise RPAServiceError("La ejecución del script excedió el tiempo límite (30 minutos)")
        
        return_code = process.returncode
        
        # Limpiar archivos temporales
        try:
            os.remove(script_file)
            if os.name == 'nt' and script_type.lower() == "python" and os.path.exists(bat_file):
                os.remove(bat_file)
        except Exception as e:
            logger.warning(f"No se pudo eliminar archivos temporales: {str(e)}")
        
        # Devolver resultado
        result = {
            "success": return_code == 0,
            "return_code": return_code,
            "stdout": stdout,
            "stderr": stderr,
            "execution_time": f"{process.returncode}s"
        }
        
        if return_code != 0:
            logger.error(f"La ejecución del script falló con código: {return_code}")
            logger.error(f"Error: {stderr}")
        else:
            logger.info("Script ejecutado exitosamente")
        
        return result
    
    except Exception as e:
        logger.error(f"Error al ejecutar script RPA: {str(e)}")
        raise RPAServiceError(f"Error al ejecutar script RPA: {str(e)}")

def execute_office_automation(app, action, file_path, params):
    """Ejecuta automatización de Office usando bibliotecas especializadas en lugar de PyAutoGUI."""
    logger.info(f"Ejecutando automatización de {app} con acción {action}")
    
    try:
        # Preparar ruta de guardado
        if not file_path:
            # Usar ruta por defecto en el escritorio
            file_path = os.path.join(os.path.expanduser('~'), 'Desktop', 
                                  f'documento_rpa_{app.lower()}.{"docx" if app.lower()=="word" else "xlsx" if app.lower()=="excel" else "pptx"}')
        
        # Asegurarse de que la ruta tenga la extensión correcta
        if app.lower() == "word" and not file_path.lower().endswith(('.doc', '.docx')):
            file_path += '.docx'
        elif app.lower() == "excel" and not file_path.lower().endswith(('.xls', '.xlsx')):
            file_path += '.xlsx'
        elif app.lower() == "powerpoint" and not file_path.lower().endswith(('.ppt', '.pptx')):
            file_path += '.pptx'
        
        logger.info(f"Usando ruta de archivo: {file_path}")
        
        # Crear directorio si no existe
        os.makedirs(os.path.dirname(os.path.abspath(file_path)), exist_ok=True)
        
        # Seleccionar la función adecuada según la aplicación
        if app.lower() == "word":
            contenido = params.get('content', 'Documento creado automáticamente')
            create_word_document(contenido, file_path)
            return {
                "success": True,
                "message": "Documento Word creado exitosamente",
                "details": f"Archivo guardado en: {file_path}"
            }
        elif app.lower() == "excel":
            datos = params.get('data', 'Datos,Valor\nEjemplo,100')
            create_excel_document(datos, file_path)
            return {
                "success": True,
                "message": "Documento Excel creado exitosamente",
                "details": f"Archivo guardado en: {file_path}"
            }
        elif app.lower() == "powerpoint":
            titulo = params.get('title', 'Presentación Automática')
            slides = params.get('slides', 'Primera Diapositiva|Contenido de la diapositiva')
            create_powerpoint_document(titulo, slides, file_path)
            return {
                "success": True,
                "message": "Presentación PowerPoint creada exitosamente",
                "details": f"Archivo guardado en: {file_path}"
            }
        else:
            raise RPAServiceError(f"Aplicación no soportada: {app}")
        
    except Exception as e:
        logger.error(f"Error en automatización de Office: {str(e)}")
        raise RPAServiceError(f"Error en automatización: {str(e)}")

def create_word_document(contenido, ruta_archivo):
    """Crea un documento Word usando python-docx."""
    try:
        # Importar la biblioteca solo cuando se necesite
        import docx
        from docx.shared import Pt
        
        # Crear un nuevo documento
        doc = docx.Document()
        
        # Agregar un título
        doc.add_heading('Documento Automatizado', 0)
        
        # Agregar contenido
        p = doc.add_paragraph(contenido)
        
        # Agregar fecha actual
        doc.add_paragraph(f'Creado automáticamente el {datetime.datetime.now().strftime("%d/%m/%Y %H:%M:%S")}')
        
        # Guardar el documento
        doc.save(ruta_archivo)
        logger.info(f"Documento Word guardado en: {ruta_archivo}")
        return True
    except ImportError:
        # Si no está instalada la biblioteca, instalarla automáticamente
        logger.warning("Biblioteca python-docx no encontrada. Instalando...")
        subprocess.check_call([sys.executable, "-m", "pip", "install", "python-docx"])
        
        # Reintentar después de instalar
        import docx
        doc = docx.Document()
        doc.add_heading('Documento Automatizado', 0)
        doc.add_paragraph(contenido)
        doc.save(ruta_archivo)
        logger.info(f"Documento Word guardado en: {ruta_archivo}")
        return True
    except Exception as e:
        logger.error(f"Error al crear documento Word: {str(e)}")
        raise RPAServiceError(f"Error al crear documento Word: {str(e)}")

def create_excel_document(datos, ruta_archivo):
    """Crea un documento Excel usando openpyxl."""
    try:
        # Importar la biblioteca solo cuando se necesite
        import openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment
        
        # Crear un nuevo libro de trabajo
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Datos Automatizados"
        
        # Procesar datos (formato CSV o líneas separadas)
        filas = datos.strip().split('\n')
        for i, fila in enumerate(filas, 1):
            valores = fila.split(',')
            for j, valor in enumerate(valores, 1):
                ws.cell(row=i, column=j, value=valor.strip())
                
                # Aplicar estilo a la primera fila (encabezados)
                if i == 1:
                    cell = ws.cell(row=i, column=j)
                    cell.font = Font(bold=True)
                    cell.fill = PatternFill(start_color="DDDDDD", end_color="DDDDDD", fill_type="solid")
        
        # Ajustar ancho de columnas automáticamente
        for col in ws.columns:
            max_length = 0
            column = col[0].column_letter
            for cell in col:
                if cell.value:
                    max_length = max(max_length, len(str(cell.value)))
            adjusted_width = (max_length + 2)
            ws.column_dimensions[column].width = adjusted_width
        
        # Guardar el libro
        wb.save(ruta_archivo)
        logger.info(f"Documento Excel guardado en: {ruta_archivo}")
        return True
    except ImportError:
        # Si no está instalada la biblioteca, instalarla automáticamente
        logger.warning("Biblioteca openpyxl no encontrada. Instalando...")
        subprocess.check_call([sys.executable, "-m", "pip", "install", "openpyxl"])
        
        # Reintentar después de instalar
        import openpyxl
        wb = openpyxl.Workbook()
        ws = wb.active
        
        filas = datos.strip().split('\n')
        for i, fila in enumerate(filas, 1):
            valores = fila.split(',')
            for j, valor in enumerate(valores, 1):
                ws.cell(row=i, column=j, value=valor.strip())
        
        wb.save(ruta_archivo)
        logger.info(f"Documento Excel guardado en: {ruta_archivo}")
        return True
    except Exception as e:
        logger.error(f"Error al crear documento Excel: {str(e)}")
        raise RPAServiceError(f"Error al crear documento Excel: {str(e)}")

def create_powerpoint_document(titulo, slides, ruta_archivo):
    """Crea una presentación PowerPoint usando python-pptx."""
    try:
        # Importar la biblioteca solo cuando se necesite
        import pptx
        from pptx.util import Inches, Pt
        
        # Crear una nueva presentación
        prs = pptx.Presentation()
        
        # Agregar slide de título
        slide_layout = prs.slide_layouts[0]  # Título y subtítulo
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = titulo
        subtitle.text = f"Creado automáticamente el {datetime.datetime.now().strftime('%d/%m/%Y')}"
        
        # Procesar slides adicionales
        if slides:
            slides_list = slides.split('|')
            if len(slides_list) >= 2:
                slide_layout = prs.slide_layouts[1]  # Título y contenido
                slide = prs.slides.add_slide(slide_layout)
                
                title = slide.shapes.title
                content = slide.placeholders[1]
                
                title.text = slides_list[0]
                content.text = slides_list[1]
        
        # Guardar la presentación
        prs.save(ruta_archivo)
        logger.info(f"Presentación PowerPoint guardada en: {ruta_archivo}")
        return True
    except ImportError:
        # Si no está instalada la biblioteca, instalarla automáticamente
        logger.warning("Biblioteca python-pptx no encontrada. Instalando...")
        subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
        
        # Reintentar después de instalar
        import pptx
        prs = pptx.Presentation()
        
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = titulo
        
        prs.save(ruta_archivo)
        logger.info(f"Presentación PowerPoint guardada en: {ruta_archivo}")
        return True
    except Exception as e:
        logger.error(f"Error al crear presentación PowerPoint: {str(e)}")
        raise RPAServiceError(f"Error al crear presentación PowerPoint: {str(e)}")